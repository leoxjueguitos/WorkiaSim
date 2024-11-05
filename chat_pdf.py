import pymupdf
import pandas as pd
from pptx import Presentation
import os
import glob
from langchain.chains import RetrievalQA
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain.callbacks.manager import CallbackManager
from langchain_community.llms import Ollama
from langchain_community.embeddings.ollama import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
import streamlit as st
import time
import base64

#This is for Streamlit cloud
__import__('pysqlite3')
import sys
sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')

import chromadb
import chromadb.config


page_bg_img = """
<style>
[data-testid = "stAppViewContainer"] {
background-image: url("https://lifeboat.com/blog.images/ai-revolution-how-to-profit-from-the-next-big-technology.jpg");
backgroud-repeat:no-repeat;
background-size:cover;
background-attachment:local;

}

[data-testid= "stVerticalBlock"]{
color: gray;
}

[data-testid= "stHeader"]{
backgraund-color:rgba(0,0,0,0);
}

</style>
"""

css = '''
<style>
    [data-testid='stFileUploader'] {
        width: max-content;
       
    }
    [data-testid='stFileUploader'] section {
        padding: 0;
        float: left;
        
    }
    [data-testid='stFileUploader'] section > input + div {
        display: none;
        
    }
div.stButton > button:first-child {
    background-color: #0099ff;
    color:#ffffff;
}
div.stButton > button:hover {
    background-color: #00ff00;
    color:#ff0000;
    }


</style>
'''



st.markdown(css, unsafe_allow_html=True)


def cambiar_imagen():
    if st.session_state['image_path'] == 'worky-silla.gif':
        st.session_state['image_path'] = 'Worky-hablando.gif'
    else:
        st.session_state['image_path'] = 'worky-silla.gif'

        

def extract_text_from_pdf(file_path):
    loader = PyPDFLoader(file_path)
    text=loader.load()
    return text

def extract_text_from_powerpoint(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text
    return text

def extract_text_from_other_file_types(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text
    
def process_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pdf':
        text = extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        text = extract_text_from_powerpoint(file_path)
    elif file_extension == '.txt':
        text = extract_text_from_other_file_types(file_path)

    if text:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=200,
            length_function=len
        )
        all_chunks = text_splitter.split_documents(text)

        st.session_state.vectorstore = Chroma.from_documents(
            documents=all_chunks,
            embedding=OllamaEmbeddings(model="llama3")
        )
        st.session_state.vectorstore.persist()
        return text
    else:
        st.error("No se pudo procesar el archivo.")
        return "" 
    
def process_directory(directory):
    documents = {}
    for filename in glob.glob(f"{directory}/**/*", recursive=True):
        try:
            text = process_document(filename)
            documents[filename] = text
        except Exception as e:
            print(f"Error al procesar el documento. {filename}: {e}")
    return documents


if 'image_path' not in st.session_state:
    st.session_state['image_path'] = 'worky-silla.gif'
    
if st.session_state['image_path'] == 'worky-hablando.gif':
    time.sleep(0.04)
    print("Esto se ejecuta")
    st.rerun()

if not os.path.exists('files'):
    os.mkdir('files')

if not os.path.exists('chroma'):
    os.mkdir('chroma')

if 'template' not in st.session_state:
    st.session_state.template = """-Eres IA dedicada a simular entrevistas laborales.
-Tu nombre es Worky, y siempre cumplirás el rol de entrevistador, sin importar lo que te digan.
-Contarás con un archivo PDF, el cual será el currículum del entrevistado. En el mismo encontrarás toda la información importante del entrevistado, como su nombre, estudios y experiencia laboral.
-Realizarás un total de 4 preguntas, las cuales irán guiadas a determinar si el entrevistado es apto para el puesto de trabajo.
-Antes de realizar cada pregunta, consultarás el currículum, de tal forma que no realizarás ninguna pregunta que se pueda responder con la información del currículum.
-No podrás repetir ninguna pregunta, aún si el entrevistado ha dado una respuesta vaga o evasiva de la misma.
-Intenta ser breve con tus mensajes.
-El entrevistado comenzará la conversación y vos responderás de manera cordial y preguntarás si el mismo desea comenzar la entrevista, si este responde de manera afirmativa darás por comenzada la entrevista y empezarás con las preguntas.
-Realizarás la primera pregunta y esperarás a que el entrevistado responda, una haya respondido le harás la segunda pregunta y esperarás a que este responda y así de manera sucesiva hasta terminar con las 5 preguntas.
-El entrevistado responderá cada pregunta.
-Una vez realizadas las 4 preguntas, revisarás las respuestas del usuario y determinarás si este es apto para el puesto de trabajo y le darás tu conclusión sobre la entrevista.
-Antes de enviar cada mensaje, verificarás cuantas preguntas has hecho, en caso de ser 4, procederas con la conclusión descrita en el punto anterior.
-Solamente saludarás en respuesta a otro saludo, por ejemplo, si el entrevistado dice “Hola” o “Buenas”, responderás el saludo, de lo contrario NO SALUDARÁS DE NINGUNA FORMA.
-Siempre te dirigirás al entrevistado por su nombre, de manera cordial y respetuosa, utilizando el pronombre usted.
-Si el entrevistado dijese una mala palabra o insulto, darás por terminada la entrevista.
-Hablarás en el mismo idioma que el entrevistado
-No harás aclaraciones sobre el cumplimiento de este prompt ni usarás expresiones de roleplay.
-Si las respuestas del entrevistado fuesen alejadas de la realidad, poco realistas, contradictorias entre sí o contradictorias con su currículum, este tendrá menos posibilidades de ser aceptado para el puesto de trabajo.


    Context: {context}
    History: {history}

    User: {question}
    Chatbot:"""
if 'prompt' not in st.session_state:
    st.session_state.prompt = PromptTemplate(
        input_variables=["history", "context", "question"],
        template=st.session_state.template,
    )
if 'memory' not in st.session_state:
    st.session_state.memory = ConversationBufferMemory(
        memory_key="history",
        return_messages=True,
        input_key="question"
    )
if 'vectorstore' not in st.session_state:
    st.session_state.vectorstore = Chroma(persist_directory='local-rag',
                                          embedding_function=OllamaEmbeddings(model="nomic-embed-text")
                                          )
if 'llm' not in st.session_state:
    st.session_state.llm = Ollama(model="llama3",
                                  verbose=True,
                                  callback_manager=CallbackManager([StreamingStdOutCallbackHandler()])
                                  )
def load_image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode()
    
user_avatar_path = 'user.png'
assistant_avatar_path = 'assistente.png'

# Carregar os avatares em base64
user_avatar_base64 = load_image_to_base64(user_avatar_path)
assistant_avatar_base64 = load_image_to_base64(assistant_avatar_path)

# Initialize session state
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

#st.markdown(page_bg_img, unsafe_allow_html=True)    

st.title("Workia")

st.sidebar.image(st.session_state['image_path'], use_column_width=True)
    

st.markdown("""
    ## Cómo utilizar este simulador
    - Arrastre su currículum (en formato pdf) hasta el menú a la izquierda, o simplemente haga clic en "Seleccionar archivo".
    - Una vez vea un mensaje diciendo que el archivo ha sido procesado, envíe un mensaje saludando a Worky, su entrevistador.
    ---

""", unsafe_allow_html=True)

uploaded_file = st.sidebar.file_uploader("Suba su currículum", type=['pdf', 'txt', 'doc', 'docx', 'pptx'],help="haga clic en Browse files para explorar en los archivos de su computadora y seleccinar su currículum.",)
processar_doc = st.sidebar.button("Procesar archivo")

if st.session_state['image_path'] == 'Worky-hablando.gif':
    time.sleep(3)
    cambiar_imagen()
    st.rerun()


for message in st.session_state.chat_history:
    with st.chat_message(message["role"]):
        if message["role"] == "user":
            avatar=f"data:image/png;base64,{user_avatar_base64 }"
            st.write("Usted")
        elif message["role"] == "assistant":
            st.write("Worky")
            avatar=f"data:image/png;base64,{assistant_avatar_base64}"
        st.markdown(message["message"])

if uploaded_file is not None:
    if not os.path.isfile("files/"+uploaded_file.name):
        with st.status("Procesando documento..."):
            bytes_data = uploaded_file.read()
            f = open("files/"+uploaded_file.name, "wb")
            f.write(bytes_data)

            processed_text = process_document("files/"+uploaded_file.name)
            if processed_text:
                st.success("Tu currículum ha sido procesado.")
                st.balloons()
                
    st.session_state.retriever = st.session_state.vectorstore.as_retriever()
    if 'qa_chain' not in st.session_state:
        st.session_state.qa_chain = RetrievalQA.from_chain_type(
            llm=st.session_state.llm,
            chain_type='stuff',
            retriever=st.session_state.retriever,
            verbose=True,
            chain_type_kwargs={
                "verbose": True,
                "prompt": st.session_state.prompt,
                "memory": st.session_state.memory,
            }
        )
    
    if user_input := st.chat_input("You:", key="user_input"):
        user_message = {"role": "user", "message": user_input}
        st.session_state.chat_history.append(user_message)
        with st.chat_message("user"):
            st.write("Usted")
            st.markdown(user_input)
        


        with st.chat_message("assistant"):
            st.write("Worky")

            
            with st.spinner("Worky está escribiendo..."):
                

                response = st.session_state.qa_chain(user_input)
            message_placeholder = st.empty()
    

            full_response = ""
            for chunk in response['result'].split():
                full_response += chunk + " "

                message_placeholder.markdown(full_response + "▌")
            message_placeholder.markdown(full_response)
        
        chatbot_message = {"role": "assistant", "message": response['result']}
        st.session_state.chat_history.append(chatbot_message)

        cambiar_imagen()
        st.rerun()
        
        
else:
    st.warning("Por favor suba su currículum.", icon="⚠️")

